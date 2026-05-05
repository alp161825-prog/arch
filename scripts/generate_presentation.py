from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PATH = "docs/筑境华章_项目展示答辩稿.pptx"

def set_bg_color(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)

def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(34, 26, 20)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12.0), Inches(0.55))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Microsoft YaHei"
        srun.font.size = Pt(15)
        srun.font.color.rgb = RGBColor(108, 88, 70)

def add_bullets(slide, items, x=0.9, y=2.0, w=11.5, h=4.8, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(45, 38, 31)
        p.space_after = Pt(8)

def add_section_card(slide, title, lines, x, y, w=3.9, h=2.0):
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(248, 242, 232)
    rect.line.color.rgb = RGBColor(214, 192, 164)
    tf = rect.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(72, 48, 29)
    for line in lines:
        sp = tf.add_paragraph()
        sp.text = line
        sp.font.name = "Microsoft YaHei"
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(70, 60, 50)

def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12.0), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(130, 110, 88)
    p.alignment = PP_ALIGN.RIGHT

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (244, 235, 220))
    add_title(s, "营造中国：古代建筑演进图谱", "中国古代建筑成就 · AI+交互信息设计参赛作品")
    add_bullets(s, ["作品定位：交互信息设计（核心），融合数据可视化与 AI 智能导览", "研究范围：民居、官署、皇宫、桥梁（1911 年以前）", "展示方式：总览主屏 → 专题探索 → 省域分析 → 知识库 → AI 助手"], y=2.2, size=20)
    add_footer(s, "团队作品展示 | 版本 2026.04")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (252, 248, 241))
    add_title(s, "项目总览", "以“一站式可探索系统”重构古代建筑知识")
    add_section_card(s, "模块 1：总览主屏", ["专题切换", "中国地图联动", "省份跳转"], 0.8, 2.0)
    add_section_card(s, "模块 2：四大专题", ["民居/官署/皇宫/桥梁", "五屏叙事", "专属交互"], 4.7, 2.0)
    add_section_card(s, "模块 3：公共分析", ["时间演进", "地域分布", "对比分析"], 8.6, 2.0)
    add_section_card(s, "模块 4：省域分析", ["高德地图", "点位联动", "图表分析"], 0.8, 4.3)
    add_section_card(s, "模块 5：知识库 + AI", ["检索/图谱/时间视图", "全局古建问答"], 4.7, 4.3, w=7.8)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (248, 243, 233))
    add_title(s, "演示路径（建议现场操作顺序）")
    add_bullets(s, ["1. 封面页点击“开始探索”，进入总览主屏，展示专题切换与地图联动。", "2. 点击省份进入省域分析页，演示地图点位 hover / click 与说明牌同步。", "3. 进入四大专题任一页，展示时间演进、结构解剖和经典案例切换。", "4. 进入知识库，演示组合检索、详情跳转、时间视图和关系图谱。", "5. 进入 AI 助手，现场提问并展示 Markdown 文档化回答效果。"], y=1.9, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (251, 246, 237))
    add_title(s, "模块一：总览主屏", "建立“全国地图 + 四类建筑 + 案例说明”的第一认知")
    add_bullets(s, ["地图交互：滚轮缩放、拖拽平移、双击复位、省份点击跳转分析页", "专题联动：左侧切换专题后，地图点位与右侧展陈说明牌同步更新", "信息密度控制：删除幕后备注文本，保留核心功能与阅读主线"], y=2.0, size=19)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (250, 245, 236))
    add_title(s, "模块二：四大专题探索", "统一五屏结构 + 专属交互，避免模板化展示")
    add_bullets(s, ["民居：地域切换 + 院落逻辑动画，突出生活方式与地域适应", "官署：空间序列 + 人物动线，展示治理秩序如何进入建筑空间", "皇宫：中轴层级高亮，解释礼制、等级与王权空间表达", "桥梁：桥型结构与受力分解，体现工程技术与交通网络"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (247, 241, 231))
    add_title(s, "模块三：省域分析页面", "从全国宏观进入区域细读")
    add_bullets(s, ["地图能力：高德底图 + 行政区高亮 + 建筑点位 + infoWindow", "交互逻辑：hover 仅提示，click 才锁定高亮并切换右侧说明牌", "图表分析：活跃度演进、城市分布、类型结构、保护与活化散点", "应用价值：支持地方文旅展示、导览路线组织与案例讲解"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (252, 247, 239))
    add_title(s, "模块四：知识库系统", "把建筑案例升级为可检索、可分析、可跳转的知识网络")
    add_bullets(s, ["搜索查询：建筑名 + 专题 + 朝代 + 省份组合检索，结果可跳详情页", "统计分析：关键词排行、热点省份矩阵、专题占比环图", "时间视图：按朝代横向浏览建筑样本，强化历史演进理解", "关系图谱：专题-朝代-建筑-省份全链路关系图，支持全屏阅读"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (249, 244, 235))
    add_title(s, "模块五：建筑详情页", "从“看图”进入“读懂建筑”")
    add_bullets(s, ["统一字段：建筑是什么、建造背景、结构特征、历史价值、观看重点", "重点建筑精细化：故宫、大明宫、二里头、赵州桥等条目深度增强", "信息互跳：可进入专题页、可进入省域页、可回到知识库继续检索"], y=2.0, size=19)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (246, 240, 230))
    add_title(s, "模块六：AI 智能助手", "面向整站的全局古建问答，而非单页附属说明")
    add_bullets(s, ["问题入口：快速提问按钮 + 自由输入，支持多轮对话", "回答能力：类型差异、朝代演进、结构逻辑、案例对比、参观理解", "输出优化：Markdown 转文档化排版（标题/列表/加粗），阅读更像 Word", "技术链路：前端请求 → 本地代理 → DeepSeek 接口"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (251, 246, 238))
    add_title(s, "数据与图像来源说明", "真实性、可解释性、可追溯性")
    add_bullets(s, ["建筑样本：围绕民居、官署、皇宫、桥梁四类，覆盖先秦至明清", "图像策略：有实景图优先替换，无实景图回退示意图，避免错配", "资料组织：专题数据、地域数据、省份数据、知识条目统一主数据源", "来源规范：官方公开资料 + 文博公开信息 + 可说明的数据整理流程"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (247, 242, 233))
    add_title(s, "创新点总结")
    add_bullets(s, ["1) 构建“时间 × 地域 × 类型 × 案例”四维知识结构", "2) 四大专题专属交互，不做模板复制", "3) 知识库与专题系统打通，支持检索、分析、图谱与详情跳转", "4) 全局 AI 问答 + 文档化输出，提升解释能力与展示体验", "5) 真实图片替换与数据字段精细化，强化内容可信度"], y=2.0, size=19)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (252, 248, 241))
    add_title(s, "10 分钟答辩时间分配（建议）")
    add_bullets(s, ["第 0-1 分钟：背景与问题定义（为什么要做这个系统）", "第 1-3 分钟：总览主屏 + 省域分析（核心交互）", "第 3-5 分钟：四大专题（结构化表达）", "第 5-7 分钟：知识库（检索、时间、图谱）", "第 7-8 分钟：AI 助手现场问答演示", "第 8-10 分钟：创新点、来源说明、应用价值与总结"], y=2.0, size=19)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (246, 240, 232))
    add_title(s, "评委高频问题预案")
    add_bullets(s, ["Q1：为什么归到“交互信息设计”？", "A：核心形态是交互系统；数据可视化是方法而非唯一呈现形式。", "Q2：数据真实性如何保证？", "A：采用公开资料来源 + 可追溯字段结构 + 重点条目精细化维护。", "Q3：AI 是否会胡说？", "A：设置了上下文约束和错误提示链路，并明确区分通识解释与实时信息。"], y=2.0, size=18)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(s, (241, 232, 218))
    add_title(s, "谢谢各位评委老师", "营造中国：古代建筑演进图谱")
    add_bullets(s, ["中国古代建筑不是孤立遗存，而是一套完整的营造文明。", "我们希望通过交互信息设计，让“看建筑”走向“理解建筑”。", "欢迎提问与指导。"], y=2.5, size=22)

    prs.save(OUT_PATH)

if __name__ == "__main__":
    create_ppt()
